import base64
import logging
import os

import httpx

from config import settings
from services import ai_service

logger = logging.getLogger(__name__)


def _validate_path(file_path: str) -> str:
    """Ensure the file path is within the uploads directory to prevent path traversal."""
    upload_dir = os.path.realpath(settings.UPLOAD_DIR)
    resolved = os.path.realpath(file_path)
    try:
        within_upload_dir = os.path.commonpath([upload_dir, resolved]) == upload_dir
    except ValueError:
        within_upload_dir = False
    if not within_upload_dir:
        raise ValueError("Access denied: path is outside the uploads directory")
    return resolved


def safe_remove_upload_file(file_path: str | None) -> bool:
    """Remove an uploaded file only after proving it lives under UPLOAD_DIR.

    Database rows can outlive old upload code and are not a safe source of
    filesystem authority. Deletion paths must therefore re-validate stored
    paths before unlinking them.
    """
    if not file_path:
        return False
    try:
        safe_path = _validate_path(file_path)
    except ValueError:
        logger.warning("Refusing to delete file outside upload directory: %s", file_path)
        return False

    if not os.path.isfile(safe_path):
        return False

    try:
        os.remove(safe_path)
    except FileNotFoundError:
        return False
    except OSError:
        logger.warning("Unable to delete upload file: %s", safe_path, exc_info=True)
        return False
    return True


def extract_text(file_path: str, file_type: str) -> str:
    """Extract text from a file based on its type."""
    safe_path = _validate_path(file_path)
    file_type = file_type.upper()

    if file_type == "PDF":
        return _extract_pdf(safe_path)
    elif file_type in ("TXT", "MD"):
        return _extract_text_file(safe_path)
    elif file_type == "PPTX":
        return _extract_pptx(safe_path)
    elif file_type == "DOCX":
        return _extract_docx(safe_path)
    else:
        raise ValueError(f"Unsupported file type for text extraction: {file_type}")


async def extract_text_async(file_path: str, file_type: str) -> str:
    """Async variant of extract_text; uses aiofiles for plain-text files."""
    safe_path = _validate_path(file_path)
    file_type_upper = file_type.upper()

    if file_type_upper in ("TXT", "MD"):
        return await _extract_text_file_async(safe_path)
    # Fall back to sync for binary formats (PDF, PPTX, DOCX use C-extensions)
    return extract_text(file_path, file_type)


def _extract_pdf(file_path: str) -> str:
    """Extract text from PDF using PyMuPDF."""
    import fitz

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    doc = fitz.open(file_path)
    text_parts = []
    for page in doc:
        text_parts.append(page.get_text())
    doc.close()
    return "\n".join(text_parts)


def _extract_text_file(file_path: str) -> str:
    """Read plain text or markdown files."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


async def _extract_text_file_async(file_path: str) -> str:
    """Async variant using aiofiles to avoid blocking the event loop."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    try:
        import aiofiles  # type: ignore
        async with aiofiles.open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return await f.read()
    except ImportError:
        # aiofiles not installed – fall back to sync read
        return _extract_text_file(file_path)


def _extract_pptx(file_path: str) -> str:
    """Extract text from PowerPoint presentations (slides + speaker notes)."""
    from pptx import Presentation

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    prs = Presentation(file_path)
    text_parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        slide_texts.append(" | ".join(row_texts))

        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        if slide_texts or notes_text:
            part = f"--- Slide {slide_num} ---\n"
            if slide_texts:
                part += "\n".join(slide_texts)
            if notes_text:
                part += f"\n[Speaker Notes]: {notes_text}"
            text_parts.append(part)

    return "\n\n".join(text_parts)


def _extract_docx(file_path: str) -> str:
    """Extract text from Word documents (paragraphs, headings, tables)."""
    from docx import Document as DocxDocument

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    doc = DocxDocument(file_path)
    text_parts: list[str] = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = paragraph.style.name if paragraph.style else ""
        if "Heading" in style_name:
            level = style_name.replace("Heading", "").strip()
            prefix = "#" * int(level) if level.isdigit() else "#"
            text_parts.append(f"{prefix} {text}")
        else:
            text_parts.append(text)

    for table in doc.tables:
        table_rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            table_rows.append(" | ".join(cells))
        if table_rows:
            text_parts.append("\n".join(table_rows))

    return "\n\n".join(text_parts)


def extract_image_text(file_path: str) -> str:
    """Extract text from an image using Groq vision model."""
    safe_path = _validate_path(file_path)

    if not os.path.exists(safe_path):
        raise FileNotFoundError(f"File not found: {safe_path}")

    if not settings.GROQ_API_KEY:
        raise ValueError("Groq API key not configured for image OCR")

    with open(safe_path, "rb") as f:
        image_data = base64.b64encode(f.read()).decode("utf-8")

    ext = os.path.splitext(safe_path)[1].lower()
    mime_map = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg"}
    mime_type = mime_map.get(ext, "image/png")

    headers = {
        "Authorization": f"Bearer {settings.GROQ_API_KEY}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": "llama-3.2-90b-vision-preview",
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": (
                            "Extract all visible text from this image. "
                            "Return a JSON object with an `extracted_text` field containing the text verbatim."
                        ),
                    },
                    {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{image_data}"}},
                ],
            }
        ],
        "max_completion_tokens": settings.LLM_MAX_COMPLETION_TOKENS,
        "temperature": settings.LLM_TEMPERATURE,
        "top_p": settings.LLM_TOP_P,
        "response_format": ai_service.JSON_RESPONSE_FORMAT if settings.LLM_JSON_MODE_ENABLED else None,
    }
    if payload["response_format"] is None:
        payload.pop("response_format")

    with httpx.Client(timeout=60.0) as client:
        response = client.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers=headers,
            json=payload,
        )
        response.raise_for_status()
        data = response.json()
        content = data["choices"][0]["message"]["content"]
        try:
            parsed = ai_service._parse_json_response(content)
            if isinstance(parsed, dict):
                return str(parsed.get("extracted_text", "")).strip()
        except (ValueError, KeyError):
            logger.warning("Failed to parse OCR JSON response, falling back to raw text")
        return content


def transcribe_audio(file_path: str) -> str:
    """Transcribe audio using Groq Whisper API."""
    safe_path = _validate_path(file_path)

    if not os.path.exists(safe_path):
        raise FileNotFoundError(f"File not found: {safe_path}")

    if not settings.GROQ_API_KEY:
        raise ValueError("Groq API key not configured for audio transcription")

    headers = {
        "Authorization": f"Bearer {settings.GROQ_API_KEY}",
    }

    with open(safe_path, "rb") as f:
        files = {"file": (os.path.basename(safe_path), f)}
        data = {"model": settings.WHISPER_MODEL}

        with httpx.Client(timeout=120.0) as client:
            response = client.post(
                "https://api.groq.com/openai/v1/audio/transcriptions",
                headers=headers,
                files=files,
                data=data,
            )
            response.raise_for_status()
            result = response.json()
            return result.get("text", "")
